from __future__ import annotations

from typing import List, Dict, Any
import os

# Existing PDF support
from pdf_processor import PDFProcessor

# Optional imports guarded at runtime to avoid import errors if not installed yet
try:
    import docx  # python-docx
except Exception:  # pragma: no cover - handled at runtime
    docx = None

try:
    from pptx import Presentation  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore

try:
    import openpyxl  # for xlsx
except Exception:  # pragma: no cover
    openpyxl = None

import csv
import io


class DocumentProcessor:
    """
    Unified document processor that extracts page-like chunks across formats:
    - .pdf -> real pages via PyPDF2
    - .docx -> treat as pseudo-pages by sectioning paragraphs into chunks
    - .pptx -> one page per slide, concatenating all text shapes
    - .xlsx/.xls -> one page per worksheet, concatenating cell values by rows
    - .csv -> one page per 500 lines (configurable)
    """

    def __init__(self):
        self.pdf = PDFProcessor()

    def extract(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(filename.lower())[1]
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        if ext == ".docx":
            return self._extract_docx(file_path)
        if ext == ".pptx":
            return self._extract_pptx(file_path)
        if ext in (".xlsx",):
            return self._extract_xlsx(file_path)
        if ext == ".csv":
            return self._extract_csv(file_path)
        raise ValueError(f"Unsupported file type: {ext}")

    # PDF
    def _extract_pdf(self, path: str) -> List[Dict[str, Any]]:
        return self.pdf.extract_pages(path)

    # DOCX
    def _extract_docx(self, path: str) -> List[Dict[str, Any]]:
        if docx is None:
            raise RuntimeError("python-docx not installed. Please add 'python-docx' to requirements.txt")
        doc = docx.Document(path)
        # Combine paragraphs and split into ~1500 char chunks as pseudo-pages
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        full_text = "\n".join(paragraphs)
        return self._chunk_text(full_text, approx_chunk_chars=1500)

    # PPTX
    def _extract_pptx(self, path: str) -> List[Dict[str, Any]]:
        if Presentation is None:
            raise RuntimeError("python-pptx not installed. Please add 'python-pptx' to requirements.txt")
        prs = Presentation(path)
        pages: List[Dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        texts.append("".join(run.text for run in p.runs) or p.text)
                elif hasattr(shape, "text"):
                    # fallback
                    texts.append(getattr(shape, "text", ""))
            content = "\n".join([t for t in texts if t])
            pages.append({
                "page_number": idx,
                "text": content or "",
                "char_count": len(content or ""),
            })
        return pages

    # XLSX/XLS
    def _extract_xlsx(self, path: str) -> List[Dict[str, Any]]:
        if openpyxl is None:
            raise RuntimeError("openpyxl not installed. Please add 'openpyxl' to requirements.txt")
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        pages: List[Dict[str, Any]] = []

        # Configuration for enriched output
        MAX_FULL_CONTENT_CHARS = 20000  # cap to keep tokens reasonable
        PREVIEW_ROWS = 10

        for idx, sheet in enumerate(wb.worksheets, start=1):
            # Collect raw rows
            raw_rows: List[list[Any]] = []
            for row in sheet.iter_rows(values_only=True):
                raw_rows.append(list(row))

            # Determine headers (Option A/B)
            headers: list[str] = []
            data_rows: List[list[Any]] = []
            if raw_rows:
                candidate = raw_rows[0]
                non_empty = [c for c in candidate if c is not None and str(c).strip() != ""]
                # Heuristic: treat first row as header if >50% non-empty and at least one string
                if len(non_empty) >= max(1, int(0.5 * len(candidate))) and any(
                    isinstance(c, str) for c in non_empty
                ):
                    headers = [str(c).strip() if c is not None else "" for c in candidate]
                    data_rows = raw_rows[1:]
                else:
                    # Generate default headers
                    max_len = max(len(r) for r in raw_rows) if raw_rows else 0
                    headers = [f"Col{i}" for i in range(1, max_len + 1)]
                    data_rows = raw_rows
            else:
                headers = []
                data_rows = []

            # Normalize row lengths to headers
            n_cols = len(headers)
            norm_rows: List[list[Any]] = []
            for r in data_rows:
                row = list(r) + [None] * (n_cols - len(r)) if n_cols > len(r) else list(r[:n_cols])
                norm_rows.append(row)

            # Build text preview and full content (Option A)
            def fmt_row(row: list[Any]) -> str:
                return "\t".join("" if v is None else str(v) for v in row)

            preview_rows = norm_rows[:PREVIEW_ROWS]
            preview_text = "\n".join(fmt_row(r) for r in preview_rows)
            full_rows_text = "\n".join(fmt_row(r) for r in norm_rows)
            full_text_truncated = full_rows_text
            truncated_note = ""
            if len(full_rows_text) > MAX_FULL_CONTENT_CHARS:
                full_text_truncated = full_rows_text[:MAX_FULL_CONTENT_CHARS] + "\n... [truncated]"
                truncated_note = f" (truncated to {MAX_FULL_CONTENT_CHARS} chars)"

            # Basic analytics (Option C without pandas)
            # Infer numeric columns
            def to_float(x: Any) -> float | None:
                if x is None:
                    return None
                try:
                    return float(x)
                except Exception:
                    return None

            col_numeric_flags: list[bool] = []
            for c in range(n_cols):
                vals = [to_float(row[c]) for row in norm_rows]
                numeric_count = sum(1 for v in vals if v is not None)
                col_numeric_flags.append(numeric_count >= max(1, int(0.6 * len(vals))) if vals else False)

            # Per-column stats
            numeric_stats: dict[str, dict[str, float | int | None]] = {}
            for c, is_num in enumerate(col_numeric_flags):
                name = headers[c] if c < len(headers) else f"Col{c+1}"
                if not is_num:
                    continue
                vals = [to_float(row[c]) for row in norm_rows if to_float(row[c]) is not None]
                if not vals:
                    numeric_stats[name] = {"count": 0, "nulls": len(norm_rows), "min": None, "max": None, "mean": None, "sum": 0}
                    continue
                s = sum(vals)
                numeric_stats[name] = {
                    "count": len(vals),
                    "nulls": len(norm_rows) - len(vals),
                    "min": min(vals),
                    "max": max(vals),
                    "mean": s / len(vals),
                    "sum": s,
                }

            # Simple groupby for first categorical column (Option C)
            cat_idx = next((i for i, f in enumerate(col_numeric_flags) if not f), None)
            groupby_summary: dict[str, dict[str, float]] = {}
            groupby_by: str | None = headers[cat_idx] if cat_idx is not None and cat_idx < len(headers) else None
            if groupby_by is not None:
                # Aggregate sums for numeric columns by category value
                aggregates: dict[str, list[float | None]] = {}
                for row in norm_rows:
                    key = "" if row[cat_idx] is None else str(row[cat_idx])
                    if key not in aggregates:
                        aggregates[key] = [0.0 if f else None for f in col_numeric_flags]
                    for c, is_num in enumerate(col_numeric_flags):
                        if not is_num:
                            continue
                        v = to_float(row[c])
                        if v is not None:
                            aggregates[key][c] = (0.0 if aggregates[key][c] is None else float(aggregates[key][c])) + v
                # Build summary for top categories by row count
                # Count occurrences
                counts: dict[str, int] = {}
                for row in norm_rows:
                    key = "" if row[cat_idx] is None else str(row[cat_idx])
                    counts[key] = counts.get(key, 0) + 1
                # Top 5 categories
                top_keys = sorted(counts.keys(), key=lambda k: counts[k], reverse=True)[:5]
                for key in top_keys:
                    per_col: dict[str, float] = {}
                    for c, is_num in enumerate(col_numeric_flags):
                        if not is_num:
                            continue
                        col_name = headers[c] if c < len(headers) else f"Col{c+1}"
                        val = aggregates.get(key, [None] * n_cols)[c]
                        if isinstance(val, (int, float)):
                            per_col[col_name] = float(val)
                    groupby_summary[key] = per_col

            # Build enriched text (Option A)
            header_line = "\t".join(headers) if headers else ""
            stats_lines = []
            for col, st in numeric_stats.items():
                stats_lines.append(
                    f"- {col}: count={st['count']}, nulls={st['nulls']}, min={st['min']}, max={st['max']}, mean={st['mean']}, sum={st['sum']}"
                )
            groupby_lines = []
            if groupby_by is not None and groupby_summary:
                groupby_lines.append(f"Grouped by '{groupby_by}' (top 5):")
                for key, agg in groupby_summary.items():
                    agg_str = ", ".join(f"{k}={v}" for k, v in agg.items())
                    groupby_lines.append(f"  - {key}: {agg_str}")

            parts: list[str] = []
            parts.append(f"Sheet: {sheet.title}")
            if headers:
                parts.append(f"Columns: {', '.join(headers)}")
            if stats_lines:
                parts.append("Numeric column stats:")
                parts.extend(stats_lines)
            if groupby_lines:
                parts.extend(groupby_lines)
            if preview_rows:
                parts.append("")
                parts.append(f"Preview (first {len(preview_rows)} rows):")
                if header_line:
                    parts.append(header_line)
                parts.append(preview_text)
            parts.append("")
            parts.append(f"Full content (tab-delimited){truncated_note}:")
            if header_line:
                parts.append(header_line)
            parts.append(full_text_truncated)

            text = "\n".join(parts)

            page_dict: Dict[str, Any] = {
                "page_number": idx,
                "text": text,
                "char_count": len(text),
            }

            # Option B: add structured fields
            page_dict["table"] = {
                "sheet": sheet.title,
                "headers": headers,
                "rows_sample": [
                    {headers[i]: ("" if r[i] is None else r[i]) for i in range(n_cols)} for r in preview_rows
                ],
                "n_rows": len(norm_rows),
                "n_cols": n_cols,
            }

            page_dict["summary"] = {
                "numeric_columns": numeric_stats,
                "groupby": {"by": groupby_by, "top": groupby_summary} if groupby_by else None,
            }

            pages.append(page_dict)

        return pages

    # CSV
    def _extract_csv(self, path: str, lines_per_page: int = 500) -> List[Dict[str, Any]]:
        pages: List[Dict[str, Any]] = []
        with open(path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            buffer: List[str] = []
            page_no = 1
            for row in reader:
                buffer.append(",".join(row))
                if len(buffer) >= lines_per_page:
                    content = "\n".join(buffer)
                    pages.append({
                        "page_number": page_no,
                        "text": content,
                        "char_count": len(content),
                    })
                    buffer = []
                    page_no += 1
            if buffer:
                content = "\n".join(buffer)
                pages.append({
                    "page_number": page_no,
                    "text": content,
                    "char_count": len(content),
                })
        return pages

    # Helpers
    def _chunk_text(self, text: str, approx_chunk_chars: int = 1500) -> List[Dict[str, Any]]:
        pages: List[Dict[str, Any]] = []
        if not text:
            return [{"page_number": 1, "text": "", "char_count": 0}]
        start = 0
        page_no = 1
        n = len(text)
        while start < n:
            end = min(n, start + approx_chunk_chars)
            # try not to cut words hard
            if end < n:
                next_break = text.rfind("\n", start, end)
                if next_break != -1 and next_break > start + 200:
                    end = next_break
            chunk = text[start:end]
            pages.append({
                "page_number": page_no,
                "text": chunk,
                "char_count": len(chunk),
            })
            page_no += 1
            start = end
        return pages
